"""
Phase 7 — Document Management System router.
Handles folders, versioning, lifecycle workflow, preview, and download.
"""
import csv
import hashlib
import io
import json
import mimetypes
import os
import time
import uuid
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Optional

import httpx
import psycopg2
from bs4 import BeautifulSoup
from fastapi import APIRouter, BackgroundTasks, Depends, File, Form, HTTPException, Query, Request, UploadFile
from fastapi.responses import HTMLResponse, Response, StreamingResponse
from html import escape as html_escape
from pypdf import PdfReader, PdfWriter
from auth import UserContext, get_user_context

KEYCLOAK_URL            = os.getenv("KEYCLOAK_URL", "http://keycloak:8080/auth")
KEYCLOAK_REALM          = os.getenv("KEYCLOAK_REALM", "ekap")
KEYCLOAK_ADMIN          = os.getenv("KEYCLOAK_ADMIN", "")
KEYCLOAK_ADMIN_PASSWORD = os.getenv("KEYCLOAK_ADMIN_PASSWORD", "")

# Large documents preview fast with just the first ~5 "pages" of content; the
# rest loads in the background in growing batches (see /preview?limit=N) so the
# frontend can keep swapping in more content without one long blocking fetch.
# A "page" means different things per format — real pages for PDF, slides for
# PPTX, sheets for XLSX, and a rough character budget for text-based formats.
PREVIEW_PAGE_LIMIT     = 5
PREVIEW_CHARS_PER_PAGE = 3000

DATABASE_URL = os.getenv("DATABASE_URL", "postgresql://ekap:changeme@postgres:5432/ekap")
UPLOAD_DIR   = Path(os.getenv("UPLOAD_DIR", "/app/uploads"))
REQUIRE_AUTH = os.getenv("REQUIRE_AUTH", "false").lower() == "true"


def safe_filename(filename: str) -> str:
    """Strip any directory/traversal components from a user-supplied filename
    before it's used to build a filesystem path — Path("../../x").name still
    returns '..' verbatim if the whole string is just '..', so that's guarded
    separately rather than relying on .name alone."""
    name = Path(filename or "").name
    if not name or name in (".", ".."):
        name = "upload"
    return name

ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".txt", ".md", ".html", ".csv", ".xlsx", ".pptx",
    ".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp",
}

# ── State transitions ─────────────────────────────────────────────────────────
# (from_state, action) → to_state
TRANSITIONS: dict[tuple[str, str], str] = {
    ("draft",    "submit"):    "review",
    ("review",   "approve"):   "approved",
    ("review",   "reject"):    "draft",
    ("approved", "publish"):   "published",
    ("published","archive"):   "archived",
    ("archived", "republish"): "published",
    ("review",   "recall"):    "draft",
    ("approved", "recall"):    "draft",
}
MANAGER_ACTIONS = {"approve", "reject", "publish", "archive", "republish"}

router = APIRouter()


# ── DB helpers ────────────────────────────────────────────────────────────────

def get_conn():
    return psycopg2.connect(DATABASE_URL)


def _require_doc(cur, document_id: str, include_deleted: bool = False) -> dict:
    extra = "" if include_deleted else "AND deleted_at IS NULL"
    cur.execute(
        f"SELECT document_id, title, owner, department, classification, status, "
        f"lifecycle_state, folder_id, current_version, file_path, file_type, deleted_at "
        f"FROM documents WHERE document_id=%s {extra}",
        (document_id,),
    )
    row = cur.fetchone()
    if not row:
        raise HTTPException(404, "Document not found.")
    cols = [d[0] for d in cur.description]
    return dict(zip(cols, row))


def _record_transition(cur, document_id: str, from_state: str, to_state: str,
                       user: UserContext, comment: str = ""):
    cur.execute(
        "INSERT INTO workflow_transitions (document_id, from_state, to_state, user_id, username, comment) "
        "VALUES (%s,%s,%s,%s,%s,%s)",
        (document_id, from_state, to_state, user.user_id, user.username, comment),
    )


# ── Branding ───────────────────────────────────────────────────────────────────
# Company/organization name shown in the admin sidebar and Employee Portal
# topbar. Read is unrestricted (any authenticated caller — both UIs need it to
# render their header); only the write is admin-gated.

@router.get("/api/branding")
async def get_branding(user: UserContext = Depends(get_user_context)):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute("SELECT company_name FROM app_branding WHERE id = 1")
            row = cur.fetchone()
    finally:
        conn.close()
    return {"company_name": row[0] if row else "EKAP"}


@router.post("/api/branding")
async def set_branding(request: Request, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    body         = await request.json()
    company_name = (body.get("company_name") or "").strip()
    if not company_name:
        raise HTTPException(400, "company_name is required.")
    if len(company_name) > 80:
        raise HTTPException(400, "company_name must be 80 characters or fewer.")

    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO app_branding (id, company_name, updated_by) VALUES (1, %s, %s) "
                "ON CONFLICT (id) DO UPDATE SET "
                "  company_name = EXCLUDED.company_name, updated_by = EXCLUDED.updated_by, updated_at = NOW()",
                (company_name, user.username),
            )
            _audit(cur, user.user_id, "BRANDING_CHANGED", None, {"company_name": company_name})
        conn.commit()
    finally:
        conn.close()
    return {"company_name": company_name}


# ── Folders ───────────────────────────────────────────────────────────────────

@router.post("/api/folders", status_code=201)
async def create_folder(
    request: Request,
    user: UserContext = Depends(get_user_context),
):
    body = await request.json()
    name = (body.get("name") or "").strip()
    if not name:
        raise HTTPException(400, "name is required.")
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    folder_id        = str(uuid.uuid4())
    parent_folder_id = body.get("parent_folder_id")
    department       = body.get("department", "")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            if parent_folder_id:
                cur.execute("SELECT 1 FROM folders WHERE folder_id=%s", (parent_folder_id,))
                if not cur.fetchone():
                    raise HTTPException(404, "Parent folder not found.")
            cur.execute(
                "INSERT INTO folders (folder_id, name, parent_folder_id, owner, department) "
                "VALUES (%s,%s,%s,%s,%s)",
                (folder_id, name, parent_folder_id, user.username, department),
            )
        conn.commit()
    finally:
        conn.close()
    return {"folder_id": folder_id, "name": name, "parent_folder_id": parent_folder_id}


@router.get("/api/folders")
async def list_folders(user: UserContext = Depends(get_user_context)):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT folder_id, name, parent_folder_id, owner, department, created_at "
                "FROM folders ORDER BY name"
            )
            cols = [d[0] for d in cur.description]
            rows = [dict(zip(cols, r)) for r in cur.fetchall()]
        # Build tree
        by_id   = {r["folder_id"]: {**r, "children": [], "created_at": str(r["created_at"])} for r in rows}
        roots   = []
        for node in by_id.values():
            pid = node["parent_folder_id"]
            if pid and pid in by_id:
                by_id[pid]["children"].append(node)
            else:
                roots.append(node)
        return {"folders": roots}
    finally:
        conn.close()


@router.get("/api/folders/{folder_id}")
async def get_folder(folder_id: str, user: UserContext = Depends(get_user_context)):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute("SELECT * FROM folders WHERE folder_id=%s", (folder_id,))
            row = cur.fetchone()
            if not row:
                raise HTTPException(404, "Folder not found.")
            cols  = [d[0] for d in cur.description]
            folder = dict(zip(cols, row))
            # Documents in this folder — same classification gate as GET /api/documents,
            # so browsing a folder can't surface titles/metadata a user isn't cleared for.
            filters, params = ["folder_id=%s", "deleted_at IS NULL"], [folder_id]
            if not user.can_manage_documents():
                placeholders = ",".join(["%s"] * len(user.accessible_classifications))
                filters.append(f"classification IN ({placeholders})")
                params.extend(list(user.accessible_classifications))
            cur.execute(
                f"SELECT document_id, title, owner, classification, status, lifecycle_state, "
                f"file_type, page_count, current_version, created_at "
                f"FROM documents WHERE {' AND '.join(filters)} ORDER BY title",
                params,
            )
            dcols = [d[0] for d in cur.description]
            docs  = [dict(zip(dcols, r)) for r in cur.fetchall()]
        folder["created_at"] = str(folder["created_at"])
        folder["updated_at"] = str(folder.get("updated_at", ""))
        return {"folder": folder, "documents": docs}
    finally:
        conn.close()


@router.put("/api/folders/{folder_id}")
async def update_folder(folder_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    body = await request.json()
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute("SELECT 1 FROM folders WHERE folder_id=%s", (folder_id,))
            if not cur.fetchone():
                raise HTTPException(404, "Folder not found.")
            fields, params = [], []
            if "name" in body:
                fields.append("name=%s"); params.append(body["name"])
            if "parent_folder_id" in body:
                fields.append("parent_folder_id=%s"); params.append(body["parent_folder_id"])
            if "department" in body:
                fields.append("department=%s"); params.append(body["department"])
            if not fields:
                raise HTTPException(400, "Nothing to update.")
            fields.append("updated_at=NOW()")
            cur.execute(f"UPDATE folders SET {','.join(fields)} WHERE folder_id=%s",
                        params + [folder_id])
        conn.commit()
    finally:
        conn.close()
    return {"status": "updated", "folder_id": folder_id}


@router.delete("/api/folders/{folder_id}", status_code=204)
async def delete_folder(folder_id: str, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute("SELECT 1 FROM folders WHERE folder_id=%s", (folder_id,))
            if not cur.fetchone():
                raise HTTPException(404, "Folder not found.")
            # Move documents to root (un-assign folder)
            cur.execute("UPDATE documents SET folder_id=NULL WHERE folder_id=%s", (folder_id,))
            cur.execute("DELETE FROM folders WHERE folder_id=%s", (folder_id,))
        conn.commit()
    finally:
        conn.close()


@router.put("/api/documents/{document_id}/folder")
async def move_document_to_folder(
    document_id: str, request: Request, user: UserContext = Depends(get_user_context)
):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    body      = await request.json()
    folder_id = body.get("folder_id")  # None = move to root
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute("SELECT 1 FROM documents WHERE document_id=%s AND deleted_at IS NULL", (document_id,))
            if not cur.fetchone():
                raise HTTPException(404, "Document not found.")
            if folder_id:
                cur.execute("SELECT 1 FROM folders WHERE folder_id=%s", (folder_id,))
                if not cur.fetchone():
                    raise HTTPException(404, "Folder not found.")
            cur.execute("UPDATE documents SET folder_id=%s, updated_at=NOW() WHERE document_id=%s",
                        (folder_id, document_id))
        conn.commit()
    finally:
        conn.close()
    return {"status": "moved", "document_id": document_id, "folder_id": folder_id}


# ── Versions ──────────────────────────────────────────────────────────────────

@router.get("/api/documents/{document_id}/versions")
async def list_versions(document_id: str, user: UserContext = Depends(get_user_context)):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            doc = _require_doc(cur, document_id)
            if doc["classification"] not in user.accessible_classifications and not user.can_manage_documents():
                raise HTTPException(403, "Access denied.")
            cur.execute(
                "SELECT version_id, version_number, file_type, file_size, uploaded_by, change_note, created_at "
                "FROM document_versions WHERE document_id=%s ORDER BY version_number DESC",
                (document_id,),
            )
            cols = [d[0] for d in cur.description]
            return {
                "document_id": document_id,
                "current_version": doc["current_version"],
                "versions": [dict(zip(cols, r)) for r in cur.fetchall()],
            }
    finally:
        conn.close()


@router.post("/api/documents/{document_id}/versions/upload", status_code=202)
async def upload_new_version(
    document_id: str,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    change_note: str = Form(""),
    user: UserContext = Depends(get_user_context),
):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    suffix = Path(file.filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise HTTPException(415, f"Unsupported file type '{suffix}'.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            doc = _require_doc(cur, document_id)
            new_ver = doc["current_version"] + 1
            dest = UPLOAD_DIR / document_id / f"v{new_ver}_{safe_filename(file.filename)}"
            dest.parent.mkdir(parents=True, exist_ok=True)
            content = await file.read()
            dest.write_bytes(content)
            file_size    = dest.stat().st_size
            content_hash = hashlib.sha256(content).hexdigest()
            cur.execute(
                "INSERT INTO document_versions (document_id, version_number, file_path, file_type, "
                "file_size, uploaded_by, change_note) VALUES (%s,%s,%s,%s,%s,%s,%s)",
                (document_id, new_ver, str(dest), suffix, file_size, user.username, change_note),
            )
            cur.execute(
                "UPDATE documents SET file_path=%s, file_type=%s, current_version=%s, content_hash=%s, "
                "status='pending', lifecycle_state='draft', updated_at=NOW() WHERE document_id=%s",
                (str(dest), suffix, new_ver, content_hash, document_id),
            )
            _record_transition(cur, document_id, doc["lifecycle_state"], "draft", user,
                               f"Version {new_ver} uploaded")
        conn.commit()
    finally:
        conn.close()

    metadata = {
        "title": doc["title"], "owner": doc["owner"],
        "department": doc["department"], "classification": doc["classification"],
    }
    background_tasks.add_task(_run_pipeline, document_id, dest, suffix, metadata)
    return {"status": "queued", "document_id": document_id, "version": new_ver}


@router.post("/api/documents/{document_id}/versions/{version_number}/restore", status_code=202)
async def restore_version(
    document_id: str,
    version_number: int,
    background_tasks: BackgroundTasks,
    user: UserContext = Depends(get_user_context),
):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            doc = _require_doc(cur, document_id)
            cur.execute(
                "SELECT file_path, file_type FROM document_versions "
                "WHERE document_id=%s AND version_number=%s",
                (document_id, version_number),
            )
            ver = cur.fetchone()
            if not ver:
                raise HTTPException(404, f"Version {version_number} not found.")
            file_path, file_type = ver
            cur.execute(
                "UPDATE documents SET file_path=%s, file_type=%s, status='pending', "
                "lifecycle_state='draft', updated_at=NOW() WHERE document_id=%s",
                (file_path, file_type, document_id),
            )
            _record_transition(cur, document_id, doc["lifecycle_state"], "draft", user,
                               f"Restored to version {version_number}")
        conn.commit()
    finally:
        conn.close()

    metadata = {
        "title": doc["title"], "owner": doc["owner"],
        "department": doc["department"], "classification": doc["classification"],
    }
    background_tasks.add_task(_run_pipeline, document_id, Path(file_path), file_type, metadata)
    return {"status": "queued", "document_id": document_id, "restored_to_version": version_number}


def _extract_version_text(path_str: str, file_type: str, version_label: str) -> str:
    """Re-run the same text extractor used at ingestion time against a stored
    version's file, so the diff compares what the AI/search actually indexed —
    not a raw byte diff, which wouldn't mean much for PDF/DOCX/XLSX/etc."""
    import main as m  # deferred: dms.py is imported by main.py, so avoid a cycle

    path = Path(path_str)
    if not path.exists():
        raise HTTPException(404, f"File for {version_label} not found on disk.")
    suffix = file_type.lower() if file_type.startswith(".") else f".{file_type}"
    try:
        chunks = m.extract_document(path, suffix)
    except Exception as exc:
        raise HTTPException(500, f"Could not extract text from {version_label}: {exc}") from exc
    return "\n\n".join(c["text"] for c in chunks)


@router.get("/api/documents/{document_id}/versions/diff")
async def diff_versions(
    document_id: str,
    from_version: int = Query(..., alias="from", ge=1),
    to_version: int = Query(..., alias="to", ge=1),
    user: UserContext = Depends(get_user_context),
):
    import difflib

    conn = get_conn()
    try:
        with conn.cursor() as cur:
            doc = _require_doc(cur, document_id)
            if doc["classification"] not in user.accessible_classifications and not user.can_manage_documents():
                raise HTTPException(403, "Access denied.")
            cur.execute(
                "SELECT version_number, file_path, file_type FROM document_versions "
                "WHERE document_id=%s AND version_number IN (%s,%s)",
                (document_id, from_version, to_version),
            )
            by_version = {r[0]: (r[1], r[2]) for r in cur.fetchall()}
    finally:
        conn.close()

    if from_version not in by_version:
        raise HTTPException(404, f"Version {from_version} not found.")
    if to_version not in by_version:
        raise HTTPException(404, f"Version {to_version} not found.")

    from_path, from_type = by_version[from_version]
    to_path, to_type     = by_version[to_version]
    from_text = _extract_version_text(from_path, from_type, f"version {from_version}")
    to_text   = _extract_version_text(to_path, to_type, f"version {to_version}")

    from_lines = from_text.splitlines()
    to_lines   = to_text.splitlines()

    # unified_diff windows context to ~3 lines either side of a change, so the
    # response stays small even for a large document — it never dumps every
    # unchanged line the way a naive full side-by-side diff would.
    raw_diff = difflib.unified_diff(from_lines, to_lines, lineterm="", n=3)
    lines, added, removed = [], 0, 0
    for line in raw_diff:
        if line.startswith(("+++", "---")):
            continue
        if line.startswith("@@"):
            lines.append({"type": "hunk", "text": line})
        elif line.startswith("+"):
            lines.append({"type": "added", "text": line[1:]})
            added += 1
        elif line.startswith("-"):
            lines.append({"type": "removed", "text": line[1:]})
            removed += 1
        else:
            lines.append({"type": "equal", "text": line[1:] if line.startswith(" ") else line})

    return {
        "document_id": document_id,
        "from_version": from_version, "to_version": to_version,
        "added_lines": added, "removed_lines": removed,
        "lines": lines,
    }


def _run_pipeline(document_id: str, file_path: Path, suffix: str, metadata: dict):
    import main as m
    import asyncio
    asyncio.run(m.process_document(document_id, file_path, suffix, metadata))


# ── Keycloak user lookup ──────────────────────────────────────────────────────
# Powers the username autocomplete on the document access-grant form, so admins
# don't have to know/type an exact user_id from memory.

_kc_admin_token: Optional[str] = None
_kc_admin_token_exp: float = 0.0


async def _get_kc_admin_token() -> str:
    global _kc_admin_token, _kc_admin_token_exp
    if _kc_admin_token and time.time() < _kc_admin_token_exp:
        return _kc_admin_token
    async with httpx.AsyncClient(timeout=10.0) as client:
        resp = await client.post(
            f"{KEYCLOAK_URL}/realms/master/protocol/openid-connect/token",
            data={
                "grant_type": "password",
                "client_id":  "admin-cli",
                "username":   KEYCLOAK_ADMIN,
                "password":   KEYCLOAK_ADMIN_PASSWORD,
            },
        )
        resp.raise_for_status()
        body = resp.json()
    _kc_admin_token     = body["access_token"]
    _kc_admin_token_exp = time.time() + body.get("expires_in", 60) - 10  # refresh a bit early
    return _kc_admin_token


@router.get("/api/keycloak/users")
async def search_keycloak_users(q: str = "", user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    if not KEYCLOAK_ADMIN or not KEYCLOAK_ADMIN_PASSWORD:
        raise HTTPException(503, "Keycloak admin credentials are not configured on this deployment.")
    try:
        token = await _get_kc_admin_token()
        async with httpx.AsyncClient(timeout=10.0) as client:
            resp = await client.get(
                f"{KEYCLOAK_URL}/admin/realms/{KEYCLOAK_REALM}/users",
                params={"search": q, "max": 10, "briefRepresentation": "true"},
                headers={"Authorization": f"Bearer {token}"},
            )
            resp.raise_for_status()
            kc_users = resp.json()
    except httpx.HTTPError as exc:
        raise HTTPException(502, f"Could not reach Keycloak: {exc}") from exc

    return {"users": [
        {"id": u.get("id"), "username": u.get("username"),
         "email": u.get("email"), "name": " ".join(filter(None, [u.get("firstName"), u.get("lastName")]))}
        for u in kc_users
    ]}


# ── Document-level permissions ──────────────────────────────────────────────
# Grants a specific user read access to a document outside their classification
# level. Enforced in retrieval-service's get_permitted_doc_ids() — this is the
# write path (grant/revoke) for that table.

def _audit(cur, user_id: str, event_type: str, resource_id: str, details: dict):
    cur.execute(
        "INSERT INTO audit_log (user_id, event_type, resource_id, details) VALUES (%s,%s,%s,%s)",
        (user_id, event_type, resource_id, json.dumps(details)),
    )


@router.get("/api/documents/{document_id}/permissions")
async def list_permissions(document_id: str, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            _require_doc(cur, document_id)
            cur.execute(
                "SELECT id, user_id, can_read, granted_by, granted_at "
                "FROM permissions WHERE document_id=%s ORDER BY granted_at DESC",
                (document_id,),
            )
            cols = [d[0] for d in cur.description]
            return {"document_id": document_id, "permissions": [dict(zip(cols, r)) for r in cur.fetchall()]}
    finally:
        conn.close()


@router.post("/api/documents/{document_id}/permissions", status_code=201)
async def grant_permission(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    body        = await request.json()
    target_user = (body.get("user_id") or "").strip()
    if not target_user:
        raise HTTPException(400, "user_id is required.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            _require_doc(cur, document_id)
            cur.execute(
                "INSERT INTO permissions (user_id, document_id, can_read, granted_by) "
                "VALUES (%s,%s,TRUE,%s) "
                "ON CONFLICT (user_id, document_id) DO UPDATE SET can_read=TRUE, granted_by=%s, granted_at=NOW()",
                (target_user, document_id, user.username, user.username),
            )
            _audit(cur, user.user_id, "PERMISSION_GRANTED", document_id,
                   {"granted_to": target_user, "granted_by": user.username})
        conn.commit()
    finally:
        conn.close()
    return {"status": "granted", "document_id": document_id, "user_id": target_user}


@router.delete("/api/documents/{document_id}/permissions/{target_user_id}", status_code=204)
async def revoke_permission(document_id: str, target_user_id: str, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            _require_doc(cur, document_id)
            cur.execute("DELETE FROM permissions WHERE document_id=%s AND user_id=%s",
                        (document_id, target_user_id))
            _audit(cur, user.user_id, "PERMISSION_REVOKED", document_id,
                   {"revoked_from": target_user_id, "revoked_by": user.username})
        conn.commit()
    finally:
        conn.close()


# ── Audit log (tamper-evident hash chain) ─────────────────────────────────────
# Every row is chained to the previous one via a Postgres trigger
# (trg_audit_log_hash_chain, see postgres/init.sql) — editing, deleting, or
# reordering any row breaks the chain from that point on. /verify recomputes
# the whole chain in SQL (same digest()/encode() the trigger uses) rather than
# in Python, so there's no risk of a JSON/text serialization mismatch producing
# a false positive.

def _audit_log_filters(
    event_type: Optional[str], user_id: Optional[str],
    start: Optional[str], end: Optional[str],
) -> tuple[list[str], list]:
    filters, params = [], []
    if event_type:
        filters.append("event_type = %s"); params.append(event_type)
    if user_id:
        filters.append("user_id = %s"); params.append(user_id)
    if start:
        filters.append("timestamp >= %s"); params.append(start)
    if end:
        filters.append("timestamp <= %s"); params.append(end)
    return filters, params


@router.get("/api/audit-log")
async def list_audit_log(
    event_type: Optional[str] = None,
    user_id: Optional[str] = None,
    start: Optional[str] = None,
    end: Optional[str] = None,
    limit: int = 50,
    offset: int = 0,
    user: UserContext = Depends(get_user_context),
):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    filters, params = _audit_log_filters(event_type, user_id, start, end)
    where = ("WHERE " + " AND ".join(filters)) if filters else ""
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                f"SELECT id, timestamp, user_id, event_type, resource_id, details, prev_hash, row_hash "
                f"FROM audit_log {where} ORDER BY id DESC LIMIT %s OFFSET %s",
                params + [limit, offset],
            )
            cols = [d[0] for d in cur.description]
            rows = [dict(zip(cols, r)) for r in cur.fetchall()]
            cur.execute(f"SELECT COUNT(*) FROM audit_log {where}", params)
            total = cur.fetchone()[0]
        return {"entries": rows, "total": total, "limit": limit, "offset": offset}
    finally:
        conn.close()


@router.get("/api/audit-log/export")
async def export_audit_log(
    event_type: Optional[str] = None,
    user_id: Optional[str] = None,
    start: Optional[str] = None,
    end: Optional[str] = None,
    user: UserContext = Depends(get_user_context),
):
    """CSV export of a date/time-ranged slice of the audit log — includes the
    hash-chain columns so the exported subset can be independently verified
    offline later. Unlike the browse endpoint, this has no pagination cap:
    it's meant to produce a complete record for a given range."""
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    filters, params = _audit_log_filters(event_type, user_id, start, end)
    where = ("WHERE " + " AND ".join(filters)) if filters else ""
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                f"SELECT id, timestamp, user_id, event_type, resource_id, details, prev_hash, row_hash "
                f"FROM audit_log {where} ORDER BY id ASC",
                params,
            )
            rows = cur.fetchall()
            _audit(cur, user.user_id, "AUDIT_LOG_EXPORTED", None,
                   {"by": user.username, "start": start, "end": end,
                    "event_type": event_type, "user_id_filter": user_id, "rows": len(rows)})
        conn.commit()
    finally:
        conn.close()

    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["id", "timestamp", "user_id", "event_type", "resource_id", "details", "prev_hash", "row_hash"])
    for r in rows:
        row = list(r)
        row[5] = json.dumps(row[5])  # details: JSONB -> JSON string for a plain CSV cell
        writer.writerow(row)

    filename = f"audit-log-{(start or 'all').replace(':', '')}-{(end or 'all').replace(':', '')}.csv"
    return Response(
        content=buf.getvalue(),
        media_type="text/csv",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/api/audit-log/verify")
async def verify_audit_log(user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute("SELECT COUNT(*) FROM audit_log")
            total = cur.fetchone()[0]
            cur.execute("""
                WITH chain_check AS (
                  SELECT
                    id,
                    prev_hash,
                    row_hash,
                    COALESCE(LAG(row_hash) OVER (ORDER BY id), '') AS expected_prev_hash,
                    encode(
                        digest(
                            COALESCE(LAG(row_hash) OVER (ORDER BY id), '') || '|' ||
                            COALESCE(user_id, '') || '|' ||
                            event_type || '|' ||
                            COALESCE(resource_id, '') || '|' ||
                            COALESCE(details::text, '{}') || '|' ||
                            timestamp::text,
                            'sha256'
                        ),
                        'hex'
                    ) AS expected_hash
                  FROM audit_log
                )
                SELECT id FROM chain_check
                WHERE row_hash IS DISTINCT FROM expected_hash
                   OR prev_hash IS DISTINCT FROM expected_prev_hash
                ORDER BY id LIMIT 1
            """)
            broken = cur.fetchone()
        return {
            "valid": broken is None,
            "checked": total,
            "first_broken_id": broken[0] if broken else None,
        }
    finally:
        conn.close()


# ── In-app notifications ──────────────────────────────────────────────────────
# user_id IS NULL rows are broadcast to anyone who can manage documents (e.g.
# "submitted for review"); note that marking a broadcast row read affects every
# manager, since read state isn't tracked per-recipient — acceptable for a
# first pass, but worth a proper junction table if this becomes noisy.

def _notification_scope(user: UserContext) -> tuple[str, list]:
    clauses, params = ["user_id=%s"], [user.username]
    if user.can_manage_documents():
        clauses.append("user_id IS NULL")
    return "(" + " OR ".join(clauses) + ")", params


@router.get("/api/notifications")
async def list_notifications(
    unread_only: bool = False, limit: int = 20, user: UserContext = Depends(get_user_context)
):
    where, params = _notification_scope(user)
    if unread_only:
        where += " AND read_at IS NULL"
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                f"SELECT notification_id, event_type, document_id, message, read_at, created_at "
                f"FROM notifications WHERE {where} ORDER BY created_at DESC LIMIT %s",
                params + [limit],
            )
            cols = [d[0] for d in cur.description]
            rows = [dict(zip(cols, r)) for r in cur.fetchall()]
            for r in rows:
                r["document_id"] = str(r["document_id"]) if r["document_id"] else None
            return {"notifications": rows}
    finally:
        conn.close()


@router.get("/api/notifications/unread-count")
async def unread_notification_count(user: UserContext = Depends(get_user_context)):
    where, params = _notification_scope(user)
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(f"SELECT COUNT(*) FROM notifications WHERE {where} AND read_at IS NULL", params)
            return {"unread": cur.fetchone()[0]}
    finally:
        conn.close()


@router.post("/api/notifications/{notification_id}/read", status_code=204)
async def mark_notification_read(notification_id: str, user: UserContext = Depends(get_user_context)):
    where, params = _notification_scope(user)
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                f"UPDATE notifications SET read_at=NOW() WHERE notification_id=%s AND {where}",
                [notification_id] + params,
            )
        conn.commit()
    finally:
        conn.close()


@router.post("/api/notifications/read-all", status_code=204)
async def mark_all_notifications_read(user: UserContext = Depends(get_user_context)):
    where, params = _notification_scope(user)
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(f"UPDATE notifications SET read_at=NOW() WHERE {where} AND read_at IS NULL", params)
        conn.commit()
    finally:
        conn.close()


# ── Retention policy ───────────────────────────────────────────────────────────
# Per-classification retention period + what happens when it elapses. The clock
# starts at publish time (see _do_transition above); run_retention_sweep() is
# called on a timer by main.py's background loop, and can also be triggered
# on demand via /api/retention/run-now.

@router.get("/api/retention-policies")
async def list_retention_policies(user: UserContext = Depends(get_user_context)):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT classification, retention_days, action, updated_at "
                "FROM retention_policies ORDER BY classification"
            )
            cols = [d[0] for d in cur.description]
            return {"policies": [dict(zip(cols, r)) for r in cur.fetchall()]}
    finally:
        conn.close()


@router.put("/api/retention-policies/{classification}")
async def update_retention_policy(
    classification: str, request: Request, user: UserContext = Depends(get_user_context)
):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    body = await request.json()
    retention_days = body.get("retention_days")
    action = body.get("action", "archive")
    if not isinstance(retention_days, int) or retention_days < 1:
        raise HTTPException(400, "retention_days must be a positive integer.")
    if action not in ("archive", "delete"):
        raise HTTPException(400, "action must be 'archive' or 'delete'.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO retention_policies (classification, retention_days, action, updated_at) "
                "VALUES (%s,%s,%s,NOW()) ON CONFLICT (classification) "
                "DO UPDATE SET retention_days=%s, action=%s, updated_at=NOW()",
                (classification, retention_days, action, retention_days, action),
            )
        conn.commit()
    finally:
        conn.close()
    return {"classification": classification, "retention_days": retention_days, "action": action}


@router.post("/api/retention/run-now")
async def trigger_retention_sweep(user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    return run_retention_sweep()


def run_retention_sweep() -> dict:
    """Archive or soft-delete documents past their retention_until, per their
    classification's policy — anything on legal hold is always skipped."""
    conn = get_conn()
    archived, deleted = 0, 0
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT document_id, title, classification, lifecycle_state FROM documents "
                "WHERE retention_until IS NOT NULL AND retention_until <= NOW() "
                "AND legal_hold = FALSE AND deleted_at IS NULL"
            )
            due = cur.fetchall()
            for doc_id, title, classification, lifecycle_state in due:
                cur.execute(
                    "SELECT action FROM retention_policies WHERE classification=%s", (classification,)
                )
                row = cur.fetchone()
                action = row[0] if row else "archive"

                if action == "delete":
                    cur.execute("UPDATE documents SET deleted_at=NOW() WHERE document_id=%s", (doc_id,))
                    _audit(cur, "system:retention", "RETENTION_DELETED", str(doc_id), {"title": title})
                    deleted += 1
                elif lifecycle_state != "archived":
                    cur.execute(
                        "UPDATE documents SET lifecycle_state='archived', updated_at=NOW() WHERE document_id=%s",
                        (doc_id,),
                    )
                    cur.execute(
                        "INSERT INTO workflow_transitions (document_id, from_state, to_state, user_id, username, comment) "
                        "VALUES (%s,%s,'archived',%s,%s,%s)",
                        (doc_id, lifecycle_state, "system:retention", "system:retention",
                         "Automatically archived — retention period elapsed"),
                    )
                    _audit(cur, "system:retention", "RETENTION_ARCHIVED", str(doc_id), {"title": title})
                    archived += 1
        conn.commit()
    finally:
        conn.close()
    return {"archived": archived, "deleted": deleted}


# ── Legal hold ─────────────────────────────────────────────────────────────────
# Overrides retention entirely while active — a held document is never touched
# by run_retention_sweep(), regardless of how far past retention_until it is.

@router.post("/api/documents/{document_id}/legal-hold")
async def set_legal_hold(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    body   = await request.json()
    reason = (body.get("reason") or "").strip()
    if not reason:
        raise HTTPException(400, "A reason is required to place a legal hold.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            _require_doc(cur, document_id)
            cur.execute(
                "UPDATE documents SET legal_hold=TRUE, legal_hold_reason=%s, "
                "legal_hold_set_by=%s, legal_hold_set_at=NOW() WHERE document_id=%s",
                (reason, user.username, document_id),
            )
            _audit(cur, user.user_id, "LEGAL_HOLD_SET", document_id, {"reason": reason, "by": user.username})
        conn.commit()
    finally:
        conn.close()
    return {"status": "held", "document_id": document_id}


@router.delete("/api/documents/{document_id}/legal-hold", status_code=204)
async def release_legal_hold(document_id: str, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            _require_doc(cur, document_id)
            cur.execute(
                "UPDATE documents SET legal_hold=FALSE, legal_hold_reason=NULL, "
                "legal_hold_set_by=NULL, legal_hold_set_at=NULL WHERE document_id=%s",
                (document_id,),
            )
            _audit(cur, user.user_id, "LEGAL_HOLD_RELEASED", document_id, {"by": user.username})
        conn.commit()
    finally:
        conn.close()


# ── Soft delete / Trash ───────────────────────────────────────────────────────

@router.get("/api/documents/trash")
async def list_trash(user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT document_id, title, owner, classification, file_type, deleted_at "
                "FROM documents WHERE deleted_at IS NOT NULL ORDER BY deleted_at DESC"
            )
            cols = [d[0] for d in cur.description]
            return {"trash": [dict(zip(cols, r)) for r in cur.fetchall()]}
    finally:
        conn.close()


@router.post("/api/documents/{document_id}/undelete")
async def undelete_document(document_id: str, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE documents SET deleted_at=NULL, updated_at=NOW() "
                "WHERE document_id=%s AND deleted_at IS NOT NULL",
                (document_id,),
            )
            if cur.rowcount == 0:
                raise HTTPException(404, "Document not found in trash.")
        conn.commit()
    finally:
        conn.close()
    return {"status": "restored", "document_id": document_id}


# ── Workflow ──────────────────────────────────────────────────────────────────

_TRANSITION_VERB = {
    "approve": "approved", "reject": "rejected", "publish": "published",
    "archive": "archived", "republish": "republished",
}


def _notify(cur, user_id: Optional[str], event_type: str, document_id: str, message: str):
    cur.execute(
        "INSERT INTO notifications (user_id, event_type, document_id, message) VALUES (%s,%s,%s,%s)",
        (user_id, event_type, document_id, message),
    )


def _do_transition(document_id: str, action: str, user: UserContext, comment: str = ""):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            doc = _require_doc(cur, document_id)
            from_state = doc["lifecycle_state"]
            to_state   = TRANSITIONS.get((from_state, action))
            if not to_state:
                raise HTTPException(
                    400,
                    f"Cannot '{action}' a document in state '{from_state}'. "
                    f"Valid actions from '{from_state}': "
                    + str([a for (s, a) in TRANSITIONS if s == from_state]),
                )
            if action in MANAGER_ACTIONS and not user.can_manage_documents():
                raise HTTPException(403, f"'{action}' requires knowledge-manager or higher.")
            if action == "submit" and doc["owner"] != user.username and not user.can_manage_documents():
                raise HTTPException(403, "Only the document owner can submit for review.")
            cur.execute(
                "UPDATE documents SET lifecycle_state=%s, updated_at=NOW() WHERE document_id=%s",
                (to_state, document_id),
            )
            _record_transition(cur, document_id, from_state, to_state, user, comment)

            # Retention clock starts (or restarts) whenever a document goes live,
            # based on its classification's current policy.
            if to_state == "published":
                cur.execute(
                    "SELECT retention_days FROM retention_policies WHERE classification=%s",
                    (doc["classification"],),
                )
                policy = cur.fetchone()
                if policy:
                    retention_until = datetime.now(timezone.utc) + timedelta(days=policy[0])
                    cur.execute(
                        "UPDATE documents SET retention_until=%s WHERE document_id=%s",
                        (retention_until, document_id),
                    )

            if action == "submit":
                _notify(cur, None, "SUBMITTED_FOR_REVIEW", document_id,
                        f'"{doc["title"]}" was submitted for review by {user.username}')
            elif action in _TRANSITION_VERB and doc["owner"] != user.username:
                msg = f'"{doc["title"]}" was {_TRANSITION_VERB[action]}'
                if comment:
                    msg += f' — "{comment}"'
                _notify(cur, doc["owner"], action.upper(), document_id, msg)
        conn.commit()
        return {"document_id": document_id, "from_state": from_state, "to_state": to_state}
    finally:
        conn.close()


@router.post("/api/documents/{document_id}/submit")
async def submit_for_review(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    return _do_transition(document_id, "submit", user, body.get("comment", ""))

@router.post("/api/documents/{document_id}/approve")
async def approve_document(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    return _do_transition(document_id, "approve", user, body.get("comment", ""))

@router.post("/api/documents/{document_id}/reject")
async def reject_document(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    return _do_transition(document_id, "reject", user, body.get("comment", ""))

@router.post("/api/documents/{document_id}/publish")
async def publish_document(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    return _do_transition(document_id, "publish", user, body.get("comment", ""))

@router.post("/api/documents/{document_id}/archive")
async def archive_document(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    return _do_transition(document_id, "archive", user, body.get("comment", ""))

@router.post("/api/documents/{document_id}/recall")
async def recall_document(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    return _do_transition(document_id, "recall", user, body.get("comment", ""))

@router.post("/api/documents/{document_id}/republish")
async def republish_document(document_id: str, request: Request, user: UserContext = Depends(get_user_context)):
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    return _do_transition(document_id, "republish", user, body.get("comment", ""))


@router.get("/api/documents/{document_id}/workflow")
async def workflow_history(document_id: str, user: UserContext = Depends(get_user_context)):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            doc = _require_doc(cur, document_id)
            if doc["classification"] not in user.accessible_classifications and not user.can_manage_documents():
                raise HTTPException(403, "Access denied.")
            cur.execute(
                "SELECT transition_id, from_state, to_state, user_id, username, comment, created_at "
                "FROM workflow_transitions WHERE document_id=%s ORDER BY created_at DESC",
                (document_id,),
            )
            cols = [d[0] for d in cur.description]
            history = [dict(zip(cols, r)) for r in cur.fetchall()]
            for h in history:
                h["created_at"] = str(h["created_at"])
        return {
            "document_id": document_id,
            "title": doc["title"],
            "current_state": doc["lifecycle_state"],
            "history": history,
        }
    finally:
        conn.close()


@router.get("/api/workflow/pending")
async def pending_reviews(user: UserContext = Depends(get_user_context)):
    """Documents currently in 'review' state, visible to knowledge-managers."""
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT document_id, title, owner, department, classification, "
                "file_type, updated_at "
                "FROM documents "
                "WHERE lifecycle_state='review' AND deleted_at IS NULL "
                "ORDER BY updated_at ASC"
            )
            cols = [d[0] for d in cur.description]
            docs = [dict(zip(cols, r)) for r in cur.fetchall()]
            for d in docs:
                d["updated_at"] = str(d["updated_at"])
        return {"pending": docs, "count": len(docs)}
    finally:
        conn.close()


# ── Preview & Download ────────────────────────────────────────────────────────

_IMAGE_TYPES  = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp", ".gif"}
_INLINE_TYPES = {".pdf", ".html", ".txt", ".csv", ".md"}


def _iter_file(path: Path, chunk_size: int = 1 << 20):
    """Stream a file in fixed-size chunks. StreamingResponse otherwise iterates
    a raw file object line-by-line (splitting on b'\\n'), which for binary files
    with many incidental newline bytes turns into hundreds of thousands of tiny
    ASGI messages and is pathologically slow."""
    with path.open("rb") as f:
        while chunk := f.read(chunk_size):
            yield chunk


def _make_scrollable(html: str) -> str:
    """Guarantees a horizontal scrollbar for wide preview content (long CSV/TXT
    lines, wide spreadsheet tables) instead of it being silently clipped at the
    iframe's edge. Applied once, after any truncation — _truncate_html_body
    above rebuilds truncated docs as bare <html><body>...</body></html>,
    dropping whatever <head><style> the original converter set, so injecting
    this per-converter wouldn't survive truncation; this runs unconditionally
    on the final HTML instead."""
    soup = BeautifulSoup(html, "html.parser")
    style_tag = soup.new_tag("style")
    style_tag.string = "html,body{margin:0} body{overflow-x:auto} pre{overflow-x:auto}"
    if soup.head:
        soup.head.append(style_tag)
    elif soup.html:
        head = soup.new_tag("head")
        head.append(style_tag)
        soup.html.insert(0, head)
    else:
        return f"<html><head>{style_tag}</head>{html}</html>"
    return str(soup)


def _truncate_html_body(html: str, char_budget: Optional[int]) -> tuple[str, int, bool]:
    """Keep whole top-level body elements until visible text reaches char_budget,
    so truncation never cuts a tag in half. char_budget=None means no truncation.
    Returns (html, total_visible_chars, was_truncated)."""
    soup       = BeautifulSoup(html, "html.parser")
    body       = soup.body if soup.body else soup
    total_len  = len(body.get_text())
    if char_budget is None or total_len <= char_budget:
        return html, total_len, False
    acc, kept = 0, []
    for child in list(body.children):
        piece = child.get_text() if hasattr(child, "get_text") else str(child)
        acc  += len(piece)
        kept.append(str(child))
        if acc >= char_budget:
            break
    return f"<html><body>{''.join(kept)}</body></html>", total_len, True


def _to_html_docx(path: Path) -> str:
    import mammoth
    with path.open("rb") as f:
        result = mammoth.convert_to_html(f)
    # Self-contained styling — this now renders inside a sandboxed iframe with
    # its own document context, so it can no longer borrow the admin portal's
    # own page CSS the way an innerHTML-inserted div could.
    return (
        "<html><head><style>body{font-family:sans-serif;max-width:860px;margin:auto;"
        "padding:24px;line-height:1.6}</style></head>"
        f"<body>{result.value}</body></html>"
    )


def _to_html_xlsx(path: Path, limit: Optional[int] = None) -> tuple[str, int, int]:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    total_sheets  = len(wb.sheetnames)
    sheet_names   = wb.sheetnames if limit is None else wb.sheetnames[:limit]
    preview_sheets = len(sheet_names)
    html = ["<html><body style='font-family:sans-serif'>"]
    for name in sheet_names:
        html.append(f"<h3>{name}</h3><table border='1' cellpadding='4' cellspacing='0' style='white-space:nowrap'>")
        for i, row in enumerate(wb[name].iter_rows(values_only=True)):
            tag = "th" if i == 0 else "td"
            html.append("<tr>" + "".join(f"<{tag}>{v if v is not None else ''}</{tag}>" for v in row) + "</tr>")
            if i > 200:
                html.append("<tr><td colspan='99'><em>…truncated…</em></td></tr>")
                break
        html.append("</table><br>")
    wb.close()
    html.append("</body></html>")
    return "\n".join(html), total_sheets, preview_sheets


def _to_html_pptx(path: Path, limit: Optional[int] = None) -> tuple[str, int, int]:
    from pptx import Presentation
    prs = Presentation(str(path))
    total_slides   = len(prs.slides)
    slides         = list(prs.slides) if limit is None else list(prs.slides)[:limit]
    preview_slides = len(slides)
    html = ["<html><body style='font-family:sans-serif'>"]
    for i, slide in enumerate(slides, 1):
        html.append(f"<div style='border:1px solid #ccc;padding:16px;margin:12px 0'>"
                    f"<strong>Slide {i}</strong><hr>")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                html.append(f"<p>{shape.text.strip()}</p>")
        html.append("</div>")
    html.append("</body></html>")
    return "\n".join(html), total_slides, preview_slides


def _to_html_md(path: Path, char_limit: Optional[int] = None) -> tuple[str, int, bool]:
    import markdown as md_lib
    text        = path.read_text(encoding="utf-8", errors="ignore")
    total_chars = len(text)
    truncated   = False
    if char_limit and total_chars > char_limit:
        cut  = text.rfind("\n\n", 0, char_limit)
        text = text[:cut if cut > 0 else char_limit]
        truncated = True
    body = md_lib.markdown(text, extensions=["tables", "fenced_code"])
    html = (
        "<html><head><style>body{font-family:sans-serif;max-width:860px;margin:auto;padding:24px}"
        "pre{background:#f5f5f5;padding:12px;overflow:auto}table{border-collapse:collapse}"
        "td,th{border:1px solid #ccc;padding:6px}</style></head>"
        f"<body>{body}</body></html>"
    )
    return html, total_chars, truncated


@router.get("/api/documents/{document_id}/preview")
async def preview_document(
    document_id: str,
    limit: Optional[int] = Query(
        None, ge=1,
        description="How many page-equivalent units to include (pages/slides/sheets/~pages of "
                     "text). Omit for the default fast preview; increase to fetch bigger batches.",
    ),
    full: bool = Query(False, description="Return the whole document, ignoring any limit."),
    user: UserContext = Depends(get_user_context),
):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            doc = _require_doc(cur, document_id)
            if doc["classification"] not in user.accessible_classifications and not user.can_manage_documents():
                raise HTTPException(403, "Access denied.")
            # Log the initial "open preview" once; the background full-document
            # fetch that follows (limit=None + full=True) is the same logical
            # view, not a separate access event.
            if limit is None and not full:
                _audit(cur, user.user_id, "DOCUMENT_VIEWED", document_id,
                       {"title": doc["title"], "classification": doc["classification"]})
        conn.commit()
    finally:
        conn.close()

    path   = Path(doc["file_path"])
    suffix = doc["file_type"].lower() if doc["file_type"].startswith(".") else f".{doc['file_type']}"

    if not path.exists():
        raise HTTPException(404, "File not found on disk.")

    # None = unlimited (the whole document); a number = that many page-equivalent
    # units. Large documents default to PREVIEW_PAGE_LIMIT; the frontend then
    # re-requests with growing `limit` values in the background for the rest.
    effective_limit = None if full else (limit if limit is not None else PREVIEW_PAGE_LIMIT)

    # PDF -> serve inline, but large PDFs get just the current batch's pages.
    if suffix == ".pdf":
        reader      = PdfReader(str(path))
        total_pages = len(reader.pages)
        eff         = total_pages if effective_limit is None else min(effective_limit, total_pages)
        headers = {
            "Content-Disposition": f"inline; filename=\"{path.name}\"",
            "X-Total-Pages": str(total_pages),
        }

        if eff >= total_pages:
            headers["X-Preview-Pages"] = str(total_pages)
            return StreamingResponse(_iter_file(path), media_type="application/pdf", headers=headers)

        # writer.append() bulk-clones a page range far faster than looping
        # add_page() one page at a time, which scales badly on large PDFs.
        writer = PdfWriter()
        writer.append(reader, pages=(0, eff))
        buf = io.BytesIO()
        writer.write(buf)
        headers["X-Preview-Pages"] = str(eff)
        # Already fully in memory — a plain Response avoids StreamingResponse's
        # slow line-by-line default iteration over the buffer.
        return Response(content=buf.getvalue(), media_type="application/pdf", headers=headers)

    if suffix in _IMAGE_TYPES:
        mime = mimetypes.types_map.get(suffix, "image/jpeg")
        return StreamingResponse(
            _iter_file(path), media_type=mime,
            headers={"Content-Disposition": f"inline; filename=\"{path.name}\""},
        )

    # Convert to HTML for in-browser preview -- large documents preview with just
    # the current batch's worth of content (see PREVIEW_* constants above); the
    # rest loads in the background via growing /preview?limit=N requests.
    char_budget = None if effective_limit is None else effective_limit * PREVIEW_CHARS_PER_PAGE
    headers: dict[str, str] = {}
    try:
        if suffix == ".docx":
            html = _to_html_docx(path)
            html, total_chars, truncated = _truncate_html_body(html, char_budget)
            total_pages = max(1, -(-total_chars // PREVIEW_CHARS_PER_PAGE))
            headers["X-Total-Pages"]   = str(total_pages)
            headers["X-Preview-Pages"] = str(min(effective_limit, total_pages)) if truncated else str(total_pages)
        elif suffix == ".xlsx":
            html, total_sheets, preview_sheets = _to_html_xlsx(path, limit=effective_limit)
            headers["X-Total-Pages"]   = str(total_sheets)
            headers["X-Preview-Pages"] = str(preview_sheets)
        elif suffix == ".pptx":
            html, total_slides, preview_slides = _to_html_pptx(path, limit=effective_limit)
            headers["X-Total-Pages"]   = str(total_slides)
            headers["X-Preview-Pages"] = str(preview_slides)
        elif suffix == ".md":
            html, total_chars, truncated = _to_html_md(path, char_limit=char_budget)
            total_pages = max(1, -(-total_chars // PREVIEW_CHARS_PER_PAGE))
            headers["X-Total-Pages"]   = str(total_pages)
            headers["X-Preview-Pages"] = str(min(effective_limit, total_pages)) if truncated else str(total_pages)
        elif suffix == ".html":
            raw = path.read_text(encoding="utf-8", errors="ignore")
            html, total_chars, truncated = _truncate_html_body(raw, char_budget)
            total_pages = max(1, -(-total_chars // PREVIEW_CHARS_PER_PAGE))
            headers["X-Total-Pages"]   = str(total_pages)
            headers["X-Preview-Pages"] = str(min(effective_limit, total_pages)) if truncated else str(total_pages)
        elif suffix in (".txt", ".csv"):
            text        = path.read_text(encoding="utf-8", errors="ignore")
            total_chars = len(text)
            truncated   = False
            if char_budget is not None and total_chars > char_budget:
                cut  = text.rfind("\n", 0, char_budget)
                text = text[:cut if cut > 0 else char_budget]
                truncated = True
            total_pages = max(1, -(-total_chars // PREVIEW_CHARS_PER_PAGE))
            headers["X-Total-Pages"]   = str(total_pages)
            headers["X-Preview-Pages"] = str(min(effective_limit, total_pages)) if truncated else str(total_pages)
            html = f"<html><body><pre>{html_escape(text)}</pre></body></html>"
        else:
            raise HTTPException(415, f"Preview not supported for '{suffix}'. Use /download instead.")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(500, f"Preview generation failed: {exc}") from exc

    return HTMLResponse(content=_make_scrollable(html), headers=headers)


@router.get("/api/documents/{document_id}/download")
async def download_document(document_id: str, user: UserContext = Depends(get_user_context)):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            doc = _require_doc(cur, document_id)
            if doc["classification"] not in user.accessible_classifications and not user.can_manage_documents():
                raise HTTPException(403, "Access denied.")
            _audit(cur, user.user_id, "DOCUMENT_DOWNLOADED", document_id,
                   {"title": doc["title"], "classification": doc["classification"]})
        conn.commit()
    finally:
        conn.close()

    path = Path(doc["file_path"])
    if not path.exists():
        raise HTTPException(404, "File not found on disk.")

    mime = mimetypes.types_map.get(doc["file_type"], "application/octet-stream")
    return StreamingResponse(
        _iter_file(path), media_type=mime,
        headers={"Content-Disposition": f"attachment; filename=\"{path.name}\""},
    )


@router.get("/api/documents/{document_id}/access-log")
async def document_access_log(document_id: str, user: UserContext = Depends(get_user_context)):
    """Who viewed/downloaded this document and when — manager-only, since it's
    sensitive activity data, distinct from the versions/workflow audit trail."""
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or higher.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            _require_doc(cur, document_id, include_deleted=True)
            cur.execute(
                "SELECT user_id, event_type, timestamp FROM audit_log "
                "WHERE resource_id=%s AND event_type IN ('DOCUMENT_VIEWED','DOCUMENT_DOWNLOADED') "
                "ORDER BY timestamp DESC LIMIT 50",
                (document_id,),
            )
            rows = cur.fetchall()
        return {
            "document_id": document_id,
            "access_log": [
                {"user_id": r[0], "event_type": r[1], "timestamp": str(r[2])}
                for r in rows
            ],
        }
    finally:
        conn.close()
