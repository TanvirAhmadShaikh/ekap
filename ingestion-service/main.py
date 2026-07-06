import asyncio
import csv
import hashlib
import json
import os
import time
import uuid
from collections import OrderedDict
from contextlib import asynccontextmanager
from pathlib import Path

import psycopg2
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from fastembed import TextEmbedding
from fastapi import BackgroundTasks, Depends, FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from prometheus_client import Counter, generate_latest, CONTENT_TYPE_LATEST
from pptx import Presentation
from pypdf import PdfReader
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance, FieldCondition, Filter, FilterSelector, MatchValue, PointStruct, VectorParams,
)
from starlette.responses import Response
import openpyxl

from auth import UserContext, get_user_context
from routers.dms import router as dms_router, run_retention_sweep

try:
    import pytesseract
    from pdf2image import convert_from_path as _pdf2img
    from PIL import Image as _PILImage
    OCR_ENABLED = True
except ImportError:
    OCR_ENABLED = False

try:
    import ocrmypdf as _ocrmypdf_lib
    OCRMYPDF_ENABLED = True
except ImportError:
    OCRMYPDF_ENABLED = False

DATABASE_URL       = os.getenv("DATABASE_URL", "postgresql://ekap:changeme@postgres:5432/ekap")
QDRANT_URL         = os.getenv("QDRANT_URL", "http://qdrant:6333")
UPLOAD_DIR         = Path(os.getenv("UPLOAD_DIR", "/app/uploads"))


def safe_filename(filename: str) -> str:
    """Strip any directory/traversal components from a user-supplied filename
    before it's used to build a filesystem path — Path("../../x").name still
    returns '..' verbatim if the whole string is just '..', so that's guarded
    separately rather than relying on .name alone."""
    name = Path(filename or "").name
    if not name or name in (".", ".."):
        name = "upload"
    return name
COLLECTION_NAME    = "ekap_chunks"
SECTIONS_COLLECTION = "ekap_sections"
EMBEDDING_MODEL    = "BAAI/bge-small-en-v1.5"
VECTOR_DIM         = 384
CHUNK_SIZE         = 1500
CHUNK_OVERLAP      = 200
VIRTUAL_SECTION_PAGES = 10       # pages per virtual section for PDFs with no headings
MAX_PDF_PAGES    = int(os.getenv("MAX_PDF_PAGES",    "500"))  # hard cap per document
OCR_TIMEOUT      = int(os.getenv("OCR_TIMEOUT",      "15"))  # seconds per page (tesseract fallback)
OCR_BUDGET_SECS  = int(os.getenv("OCR_BUDGET_SECS",  "90"))  # total OCR wall-time per doc (tesseract fallback)
OCR_JOBS         = int(os.getenv("OCR_JOBS",          "4"))   # parallel ocrmypdf workers
# Small on purpose: onnxruntime's CPU arena grows to fit the largest batch it has
# seen and never shrinks, so a big EMBED_BATCH here permanently inflates the
# embedder's RSS for the life of the process, not just for this one document.
EMBED_BATCH      = int(os.getenv("EMBED_BATCH",       "32"))
ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".txt", ".md", ".html", ".csv", ".xlsx", ".pptx",
    ".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".webp",
}
REQUIRE_AUTH = os.getenv("REQUIRE_AUTH", "false").lower() == "true"
RETENTION_SWEEP_INTERVAL_SECONDS = int(os.getenv("RETENTION_SWEEP_INTERVAL_SECONDS", "3600"))

qdrant: QdrantClient = None
embedder: TextEmbedding = None

ingestion_counter = Counter("ekap_ingestions_total", "Documents ingested", ["status"])


async def _retention_loop():
    """Periodically archive/delete documents whose retention period elapsed.
    Nothing else in this stack runs scheduled jobs, so this is a plain
    in-process asyncio loop rather than a separate cron container."""
    while True:
        try:
            result = await asyncio.to_thread(run_retention_sweep)
            if result["archived"] or result["deleted"]:
                print(f"[retention] archived={result['archived']} deleted={result['deleted']}", flush=True)
        except Exception as exc:
            print(f"[retention] sweep failed: {exc}", flush=True)
        await asyncio.sleep(RETENTION_SWEEP_INTERVAL_SECONDS)


@asynccontextmanager
async def lifespan(app: FastAPI):
    global qdrant, embedder
    # Reset any documents that were mid-processing when the container last died
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE documents SET status='failed', processing_stage=NULL "
                "WHERE status='processing'"
            )
            if cur.rowcount:
                print(f"[startup] Reset {cur.rowcount} incomplete processing job(s) to 'failed'", flush=True)
        conn.commit()
    finally:
        conn.close()

    qdrant = QdrantClient(QDRANT_URL)
    embedder = TextEmbedding(EMBEDDING_MODEL)
    vec_params = VectorParams(size=VECTOR_DIM, distance=Distance.COSINE)
    for col in (COLLECTION_NAME, SECTIONS_COLLECTION):
        if not qdrant.collection_exists(col):
            qdrant.create_collection(collection_name=col, vectors_config=vec_params)
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    retention_task = asyncio.create_task(_retention_loop())
    yield
    retention_task.cancel()


app = FastAPI(title="EKAP Ingestion Service", version="0.7.0", lifespan=lifespan)
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
app.include_router(dms_router)


# ── Helpers ────────────────────────────────────────────────────────────────────

def get_conn():
    return psycopg2.connect(DATABASE_URL)


def _audit_sync(user_id: str, event_type: str, resource_id: str | None, details: dict):
    try:
        conn = get_conn()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO audit_log (user_id, event_type, resource_id, details) VALUES (%s,%s,%s,%s)",
                    (user_id, event_type, resource_id, json.dumps(details)),
                )
            conn.commit()
        finally:
            conn.close()
    except Exception:
        pass


async def audit(user_id: str, event_type: str, resource_id: str | None = None, details: dict | None = None):
    asyncio.create_task(asyncio.to_thread(_audit_sync, user_id, event_type, resource_id, details or {}))


def make_chunks(text: str, page_number: int, section: str | None = None) -> list[dict]:
    text = text.strip()
    if not text:
        return []
    result, start = [], 0
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        piece = text[start:end].strip()
        if piece:
            result.append({"text": piece, "page_number": page_number, "section": section})
        if end >= len(text):
            break
        start = end - CHUNK_OVERLAP
    return result


# ── Auto-OCR helpers ───────────────────────────────────────────────────────────

def _needs_ocr(path: Path, sample_count: int = 6) -> bool:
    """Return True if most sampled pages have < 50 chars (i.e. the PDF is scanned)."""
    try:
        reader = PdfReader(str(path))
        n = len(reader.pages)
        if n == 0:
            return False
        # Sample pages spread across the document (beginning, middle, end)
        fractions = [0.0, 0.2, 0.4, 0.6, 0.8, 1.0]
        indices = list(dict.fromkeys(min(int(n * f), n - 1) for f in fractions))
        scanned = sum(1 for i in indices if len(reader.pages[i].extract_text() or "") < 50)
        return scanned > len(indices) / 2
    except Exception:
        return False


def _ocrmypdf_convert(path: Path) -> Path:
    """
    Run ocrmypdf on a scanned PDF to embed an invisible text layer.
    --skip-text means pages that already have text are left untouched.
    Returns the path to the OCR'd PDF; falls back to original if conversion fails.
    """
    ocr_path = path.parent / f"{path.stem}_ocr{path.suffix}"
    # Re-use an existing OCR output so crash-and-retry doesn't re-run the full job
    if ocr_path.exists():
        print(f"[ocrmypdf] Re-using cached OCR output: {ocr_path}", flush=True)
        return ocr_path
    try:
        _ocrmypdf_lib.ocr(
            path, ocr_path,     # positional — first arg renamed in v17
            skip_text=True,     # don't re-OCR pages that already have embedded text
            jobs=OCR_JOBS,      # parallel tesseract workers
            progress_bar=False,
            optimize=0,         # skip image compression — we only need the text layer
        )
        if ocr_path.exists():
            print(f"[ocrmypdf] Created text-searchable PDF: {ocr_path}", flush=True)
            return ocr_path
    except Exception as exc:
        print(f"[ocrmypdf] Pre-processing failed, falling back to per-page OCR: {exc}", flush=True)
        if ocr_path.exists():
            ocr_path.unlink(missing_ok=True)
    return path


# ── Format extractors ──────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> list[dict]:
    reader     = PdfReader(str(path))
    total      = len(reader.pages)
    limit      = min(total, MAX_PDF_PAGES)
    ocr_spent  = 0.0   # running total of seconds spent on OCR across all pages
    chunks     = []

    for page_num, page in enumerate(reader.pages, start=1):
        if page_num > limit:
            chunks.extend(make_chunks(
                f"[Document truncated: processed {limit} of {total} pages. "
                f"Increase MAX_PDF_PAGES env var to index more.]",
                page_num,
            ))
            break
        text = page.extract_text() or ""
        if len(text.strip()) < 50 and OCR_ENABLED and ocr_spent < OCR_BUDGET_SECS:
            t0 = time.monotonic()
            try:
                imgs = _pdf2img(str(path), first_page=page_num, last_page=page_num, dpi=100)
                if imgs:
                    text = pytesseract.image_to_string(imgs[0], timeout=OCR_TIMEOUT)
            except Exception:
                pass  # timeout or conversion error — just use whatever text we have
            ocr_spent += time.monotonic() - t0
        chunks.extend(make_chunks(text, page_num))

    return chunks


def extract_image(path: Path) -> list[dict]:
    if not OCR_ENABLED:
        return []
    try:
        text = pytesseract.image_to_string(_PILImage.open(str(path)))
        return make_chunks(text, page_number=1)
    except Exception:
        return []


def extract_docx(path: Path) -> list[dict]:
    doc = DocxDocument(str(path))
    chunks, current_section, para_count = [], None, 0
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            current_section = para.text.strip() or current_section
            continue
        text = para.text.strip()
        if not text:
            continue
        para_count += 1
        chunks.extend(make_chunks(text, (para_count - 1) // 30 + 1, current_section))
    return chunks


def extract_plain_text(path: Path) -> list[dict]:
    lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    chunks = []
    for i in range(0, len(lines), 50):
        block = "\n".join(lines[i : i + 50])
        chunks.extend(make_chunks(block, i // 50 + 1))
    return chunks


def extract_markdown(path: Path) -> list[dict]:
    lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    chunks, current_section, buffer, page_num = [], None, [], 1
    for line in lines:
        if line.startswith("#"):
            if buffer:
                chunks.extend(make_chunks("\n".join(buffer), page_num, current_section))
                buffer = []
                page_num += 1
            current_section = line.lstrip("#").strip()
        else:
            buffer.append(line)
    if buffer:
        chunks.extend(make_chunks("\n".join(buffer), page_num, current_section))
    return chunks


def extract_html(path: Path) -> list[dict]:
    soup = BeautifulSoup(path.read_bytes(), "lxml")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()
    chunks, current_section, page_num = [], None, 1
    for elem in soup.find_all(["h1", "h2", "h3", "p", "li", "td", "th"]):
        if elem.name in ("h1", "h2", "h3"):
            current_section = elem.get_text(strip=True)
            page_num += 1
        else:
            text = elem.get_text(separator=" ", strip=True)
            chunks.extend(make_chunks(text, page_num, current_section))
    return chunks


def extract_csv(path: Path) -> list[dict]:
    rows = []
    with path.open(encoding="utf-8", errors="ignore", newline="") as f:
        for row in csv.reader(f):
            line = ", ".join(cell for cell in row if cell.strip())
            if line:
                rows.append(line)
    return make_chunks("\n".join(rows), page_number=1)


def extract_excel(path: Path) -> list[dict]:
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    chunks = []
    for sheet_num, name in enumerate(wb.sheetnames, start=1):
        rows = []
        for row in wb[name].iter_rows(values_only=True):
            line = ", ".join(str(v) for v in row if v is not None)
            if line:
                rows.append(line)
        if rows:
            chunks.extend(make_chunks(f"Sheet: {name}\n" + "\n".join(rows), sheet_num, name))
    wb.close()
    return chunks


def extract_pptx(path: Path) -> list[dict]:
    prs = Presentation(str(path))
    chunks = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts, title = [], None
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text.strip():
                continue
            if title is None:
                title = shape.text.strip()
            else:
                texts.append(shape.text.strip())
        slide_text = "\n".join(filter(None, [title] + texts))
        chunks.extend(make_chunks(slide_text, slide_num, title))
    return chunks


def extract_document(path: Path, suffix: str) -> list[dict]:
    _img = extract_image
    extractors = {
        ".pdf": extract_pdf, ".docx": extract_docx, ".txt": extract_plain_text,
        ".md":  extract_markdown, ".html": extract_html, ".csv": extract_csv,
        ".xlsx": extract_excel, ".pptx": extract_pptx,
        ".jpg": _img, ".jpeg": _img, ".png": _img,
        ".tiff": _img, ".tif": _img, ".bmp": _img, ".webp": _img,
    }
    fn = extractors.get(suffix)
    if fn is None:
        raise ValueError(f"No extractor for '{suffix}'")
    return fn(path)


# ── Section grouping (Phase 4) ─────────────────────────────────────────────────

def group_into_sections(chunks: list[dict]) -> list[dict]:
    """
    Group chunks by section label. Chunks without explicit sections (e.g. PDFs)
    are grouped into virtual sections of VIRTUAL_SECTION_PAGES pages each.
    Returns list of section dicts with summary (extractive) for embedding.
    """
    section_map: OrderedDict[str, dict] = OrderedDict()

    for i, chunk in enumerate(chunks):
        key = chunk.get("section")
        if not key:
            p = chunk.get("page_number", 1)
            lo = ((p - 1) // VIRTUAL_SECTION_PAGES) * VIRTUAL_SECTION_PAGES + 1
            key = f"Pages {lo}–{lo + VIRTUAL_SECTION_PAGES - 1}"

        if key not in section_map:
            section_map[key] = {
                "section_id":    str(uuid.uuid4()),
                "section_title": key,
                "start_idx":     i,
                "end_idx":       i,
                "texts":         [],
            }
        section_map[key]["end_idx"] = i
        section_map[key]["texts"].append(chunk["text"])

    sections = []
    for title, s in section_map.items():
        preview = " ".join(s["texts"])[:1000]
        sections.append({
            "section_id":       s["section_id"],
            "section_title":    title,
            "start_idx":        s["start_idx"],
            "end_idx":          s["end_idx"],
            "summary":          f"{title}. {preview}",
            "qdrant_point_id":  None,
        })
    return sections


# ── Ingestion pipeline ─────────────────────────────────────────────────────────

def _set_stage(conn, document_id: str, stage: str) -> None:
    with conn.cursor() as cur:
        # Abort if the admin cancelled this job between stages
        cur.execute("SELECT status FROM documents WHERE document_id=%s", (document_id,))
        row = cur.fetchone()
        if row and row[0] == "failed":
            raise RuntimeError(f"Processing of {document_id} was cancelled.")
        cur.execute(
            "UPDATE documents SET processing_stage=%s, updated_at=NOW() WHERE document_id=%s",
            (stage, document_id),
        )
    conn.commit()


def _purge_document_vectors(document_id: str) -> None:
    """Delete all existing vectors and DB records for this document before reindexing."""
    try:
        qdrant.delete(collection_name=COLLECTION_NAME,
                      points_selector=FilterSelector(filter=Filter(must=[
                          FieldCondition(key="document_id", match=MatchValue(value=document_id))])))
        qdrant.delete(collection_name=SECTIONS_COLLECTION,
                      points_selector=FilterSelector(filter=Filter(must=[
                          FieldCondition(key="document_id", match=MatchValue(value=document_id))])))
    except Exception as exc:
        print(f"[reindex] Qdrant purge warning for {document_id}: {exc}", flush=True)
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute("DELETE FROM chunks   WHERE document_id=%s", (document_id,))
            cur.execute("DELETE FROM sections WHERE document_id=%s", (document_id,))
        conn.commit()
    finally:
        conn.close()


def _process_document(document_id: str, file_path: Path, suffix: str, metadata: dict) -> None:
    conn = get_conn()
    try:
        # Purge any previous index data so reindexing replaces rather than appends
        _purge_document_vectors(document_id)

        with conn.cursor() as cur:
            cur.execute(
                "UPDATE documents SET status='processing', processing_stage='extracting', "
                "processing_started_at=NOW(), updated_at=NOW() WHERE document_id=%s",
                (document_id,),
            )
        conn.commit()

        # ── 0. Auto-OCR: convert scanned PDFs so pypdf can extract their text ─
        if suffix == ".pdf" and OCRMYPDF_ENABLED and _needs_ocr(file_path):
            _set_stage(conn, document_id, "ocr_converting")
            file_path = _ocrmypdf_convert(file_path)

        # ── 1. Extract ────────────────────────────────────────────────────────
        _set_stage(conn, document_id, "extracting")
        chunks = extract_document(file_path, suffix)
        if not chunks:
            raise ValueError("No text could be extracted from the document.")

        # ── 2. Group into sections ────────────────────────────────────────────
        _set_stage(conn, document_id, "sectioning")
        sections = group_into_sections(chunks)

        # ── 3. Embed section summaries ────────────────────────────────────────
        _set_stage(conn, document_id, "embedding_sections")
        section_summaries = [s["summary"] for s in sections]
        section_embeddings = list(embedder.embed(section_summaries))

        # ── 4. Upsert sections → Qdrant + PostgreSQL ──────────────────────────
        _set_stage(conn, document_id, "indexing_sections")
        section_points = []
        for s, emb in zip(sections, section_embeddings):
            s["qdrant_point_id"] = str(uuid.uuid4())
            section_points.append(PointStruct(
                id=s["qdrant_point_id"],
                vector=emb.tolist(),
                payload={
                    "section_id":     s["section_id"],
                    "document_id":    document_id,
                    "document_title": metadata["title"],
                    "section_title":  s["section_title"],
                    "summary":        s["summary"],
                    "classification": metadata["classification"],
                    "department":     metadata.get("department", ""),
                },
            ))

        for batch in range(0, len(section_points), 100):
            qdrant.upsert(collection_name=SECTIONS_COLLECTION, points=section_points[batch:batch+100])

        with conn.cursor() as cur:
            cur.executemany(
                """INSERT INTO sections
                   (section_id, document_id, section_title, summary, start_chunk_idx, end_chunk_idx, qdrant_point_id)
                   VALUES (%s, %s, %s, %s, %s, %s, %s)""",
                [(s["section_id"], document_id, s["section_title"], s["summary"],
                  s["start_idx"], s["end_idx"], s["qdrant_point_id"])
                 for s in sections],
            )
        conn.commit()

        # ── 5. Build chunk → section_id map ──────────────────────────────────
        chunk_to_section: dict[int, str] = {}
        for s in sections:
            for idx in range(s["start_idx"], s["end_idx"] + 1):
                chunk_to_section[idx] = s["section_id"]

        # ── 6. Embed chunks (slowest step — scales with document size) ────────
        _set_stage(conn, document_id, "embedding_chunks")
        texts = [c["text"] for c in chunks]
        embeddings = []
        for i in range(0, len(texts), EMBED_BATCH):
            embeddings.extend(embedder.embed(texts[i:i + EMBED_BATCH]))

        # ── 7. Upsert chunks → Qdrant ─────────────────────────────────────────
        _set_stage(conn, document_id, "indexing_chunks")
        points, db_chunks = [], []
        for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
            point_id   = str(uuid.uuid4())
            chunk_id   = str(uuid.uuid4())
            section_id = chunk_to_section.get(i)

            points.append(PointStruct(
                id=point_id,
                vector=emb.tolist(),
                payload={
                    "document_id":    document_id,
                    "chunk_id":       chunk_id,
                    "section_id":     section_id,
                    "document_title": metadata["title"],
                    "page_number":    chunk["page_number"],
                    "section":        chunk.get("section"),
                    "classification": metadata["classification"],
                    "department":     metadata.get("department", ""),
                    "owner":          metadata["owner"],
                    "chunk_index":    i,
                    "content":        chunk["text"],
                },
            ))
            db_chunks.append((
                chunk_id, document_id, section_id, i, chunk["text"],
                point_id, chunk["page_number"], chunk.get("section"),
            ))

        for batch in range(0, len(points), 100):
            qdrant.upsert(collection_name=COLLECTION_NAME, points=points[batch:batch+100])

        # ── 8. Write chunks + finalise ────────────────────────────────────────
        _set_stage(conn, document_id, "finalizing")
        with conn.cursor() as cur:
            cur.executemany(
                """INSERT INTO chunks
                   (chunk_id, document_id, section_id, chunk_index, content,
                    qdrant_point_id, page_number, section)
                   VALUES (%s, %s, %s, %s, %s, %s, %s, %s)""",
                db_chunks,
            )
            page_count = max(c["page_number"] for c in chunks)
            cur.execute(
                "UPDATE documents SET status='completed', processing_stage=NULL, "
                "page_count=%s, updated_at=NOW() WHERE document_id=%s",
                (page_count, document_id),
            )
            cur.execute(
                "INSERT INTO document_versions "
                "(document_id, version_number, file_path, file_type, file_size, uploaded_by) "
                "SELECT %s, current_version, file_path, file_type, "
                "       (SELECT octet_length(content::text) FROM chunks WHERE document_id=%s LIMIT 1), "
                "       owner "
                "FROM documents WHERE document_id=%s "
                "ON CONFLICT (document_id, version_number) DO NOTHING",
                (document_id, document_id, document_id),
            )
        conn.commit()
        ingestion_counter.labels(status="completed").inc()

    except Exception:
        conn.rollback()
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE documents SET status='failed', processing_stage=NULL, updated_at=NOW() "
                "WHERE document_id=%s",
                (document_id,),
            )
        conn.commit()
        ingestion_counter.labels(status="failed").inc()
        raise
    finally:
        conn.close()


async def process_document(document_id: str, file_path: Path, suffix: str, metadata: dict) -> None:
    await asyncio.to_thread(_process_document, document_id, file_path, suffix, metadata)


# ── Endpoints ──────────────────────────────────────────────────────────────────

@app.get("/health")
def health():
    return {"status": "ok", "service": "ingestion-service", "version": "0.7.0"}


@app.get("/metrics")
def metrics():
    return Response(generate_latest(), media_type=CONTENT_TYPE_LATEST)


@app.post("/api/documents/upload", status_code=202)
async def upload_document(
    request: Request,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    title: str = Form(...),
    owner: str = Form(...),
    department: str = Form(""),
    classification: str = Form("Internal"),
    tags: str = Form(""),
    folder_id: str = Form(None),
    force: bool = Form(False),
    user: UserContext = Depends(get_user_context),
):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager, administrator, or system-administrator role.")

    suffix = Path(file.filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise HTTPException(415, f"Unsupported file type '{suffix}'. Allowed: {sorted(ALLOWED_EXTENSIONS)}")

    content      = await file.read()
    content_hash = hashlib.sha256(content).hexdigest()

    # Duplicate detection: same file re-uploaded, or a document with this exact
    # title already exists. Both are surfaced as a 409 so the caller can choose
    # to upload as a new version of the existing document instead, or proceed
    # anyway with force=true — never silently blocked.
    if not force:
        conn = get_conn()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "SELECT document_id, title FROM documents "
                    "WHERE deleted_at IS NULL AND content_hash=%s LIMIT 1",
                    (content_hash,),
                )
                match = cur.fetchone()
                if match:
                    raise HTTPException(409, detail={
                        "conflict": "content",
                        "existing_document_id": str(match[0]),
                        "existing_title": match[1],
                        "message": f'This exact file already exists as "{match[1]}".',
                    })
                cur.execute(
                    "SELECT document_id, title FROM documents "
                    "WHERE deleted_at IS NULL AND LOWER(title)=LOWER(%s) LIMIT 1",
                    (title,),
                )
                match = cur.fetchone()
                if match:
                    raise HTTPException(409, detail={
                        "conflict": "title",
                        "existing_document_id": str(match[0]),
                        "existing_title": match[1],
                        "message": f'A document titled "{match[1]}" already exists.',
                    })
        finally:
            conn.close()

    document_id = str(uuid.uuid4())
    dest = UPLOAD_DIR / document_id / safe_filename(file.filename)
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(content)

    # In dev mode (no auth), skip approval workflow and publish immediately
    initial_lifecycle = "published" if not REQUIRE_AUTH else "draft"

    tag_list = [t.strip() for t in tags.split(",") if t.strip()]
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            if folder_id:
                cur.execute("SELECT 1 FROM folders WHERE folder_id=%s", (folder_id,))
                if not cur.fetchone():
                    raise HTTPException(404, f"Folder '{folder_id}' not found.")
            cur.execute(
                """INSERT INTO documents
                   (document_id, title, owner, department, classification, tags, status,
                    file_path, file_type, folder_id, lifecycle_state, content_hash)
                   VALUES (%s, %s, %s, %s, %s, %s, 'pending', %s, %s, %s, %s, %s)""",
                (document_id, title, owner, department, classification, tag_list,
                 str(dest), suffix, folder_id, initial_lifecycle, content_hash),
            )
        conn.commit()
    finally:
        conn.close()

    metadata = {"title": title, "owner": owner, "department": department, "classification": classification}
    background_tasks.add_task(process_document, document_id, dest, suffix, metadata)
    ingestion_counter.labels(status="queued").inc()

    await audit(user.user_id, "DOCUMENT_UPLOAD", resource_id=document_id,
                details={"title": title, "classification": classification, "filename": file.filename,
                         "folder_id": folder_id, "lifecycle_state": initial_lifecycle,
                         "username": user.username})

    return {
        "document_id": document_id, "status": "queued",
        "filename": file.filename, "format": suffix,
        "lifecycle_state": initial_lifecycle, "folder_id": folder_id,
    }


@app.get("/api/documents")
async def list_documents(
    q: str = None,
    department: str = None,
    classification: str = None,
    lifecycle_state: str = None,
    folder_id: str = None,
    include_deleted: bool = False,
    limit: int = 50,
    offset: int = 0,
    user: UserContext = Depends(get_user_context),
):
    conn = get_conn()
    try:
        filters, params = [], []
        # Always exclude soft-deleted unless manager requests them explicitly
        if not (include_deleted and user.can_manage_documents()):
            filters.append("deleted_at IS NULL")
        if q:
            # Real corpus search: title/owner match, or full-text match against
            # any chunk's content (same plainto_tsquery + content_fts index the
            # AI chat's BM25 retrieval already uses) — not just a substring
            # filter over whatever page of results happens to be loaded.
            filters.append(
                "(title ILIKE %s OR owner ILIKE %s OR EXISTS ("
                "SELECT 1 FROM chunks c WHERE c.document_id = documents.document_id "
                "AND c.content_fts @@ plainto_tsquery('english', %s)))"
            )
            like_pattern = f"%{q}%"
            params.extend([like_pattern, like_pattern, q])
        if department:
            filters.append("department = %s"); params.append(department)
        if classification:
            filters.append("classification = %s"); params.append(classification)
        if lifecycle_state:
            filters.append("lifecycle_state = %s"); params.append(lifecycle_state)
        if folder_id:
            filters.append("folder_id = %s"); params.append(folder_id)
        # Non-managers only see documents within their classification access
        if not user.can_manage_documents():
            placeholders = ",".join(["%s"] * len(user.accessible_classifications))
            filters.append(f"classification IN ({placeholders})")
            params.extend(list(user.accessible_classifications))
        where = ("WHERE " + " AND ".join(filters)) if filters else ""
        with conn.cursor() as cur:
            cur.execute(
                f"SELECT document_id, title, owner, department, classification, status, "
                f"file_type, page_count, created_at, lifecycle_state, folder_id, deleted_at, "
                f"processing_stage, processing_started_at, current_version, "
                f"retention_until, legal_hold, legal_hold_reason "
                f"FROM documents {where} ORDER BY created_at DESC LIMIT %s OFFSET %s",
                params + [limit, offset],
            )
            rows = cur.fetchall()
            cur.execute(f"SELECT COUNT(*) FROM documents {where}", params)
            total = cur.fetchone()[0]
        return {
            "documents": [
                {"document_id": str(r[0]), "title": r[1], "owner": r[2], "department": r[3],
                 "classification": r[4], "status": r[5], "file_type": r[6],
                 "page_count": r[7], "created_at": str(r[8]), "lifecycle_state": r[9],
                 "folder_id": str(r[10]) if r[10] else None,
                 "deleted_at": str(r[11]) if r[11] else None,
                 "processing_stage": r[12],
                 "processing_started_at": str(r[13]) if r[13] else None,
                 "current_version": r[14],
                 "retention_until": str(r[15]) if r[15] else None,
                 "legal_hold": r[16], "legal_hold_reason": r[17]}
                for r in rows
            ],
            "total": total, "limit": limit, "offset": offset,
        }
    finally:
        conn.close()


@app.get("/api/documents/{document_id}")
async def get_document(document_id: str, user: UserContext = Depends(get_user_context)):
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT document_id, title, owner, department, classification, status, "
                "page_count, file_type, tags, created_at, updated_at, "
                "lifecycle_state, folder_id, current_version, deleted_at, "
                "processing_stage, processing_started_at "
                "FROM documents WHERE document_id = %s AND deleted_at IS NULL",
                (document_id,),
            )
            row = cur.fetchone()
            if not row:
                raise HTTPException(404, "Document not found.")
            cols = [d[0] for d in cur.description]
            doc = {c: (str(v) if v is not None else None) for c, v in zip(cols, row)}
        # Check classification access
        if doc["classification"] not in user.accessible_classifications and not user.can_manage_documents():
            raise HTTPException(403, "You do not have access to this document.")
        return doc
    finally:
        conn.close()


@app.delete("/api/documents/{document_id}", status_code=204)
async def delete_document(
    document_id: str,
    user: UserContext = Depends(get_user_context),
):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager, administrator, or system-administrator role.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT document_id, title FROM documents WHERE document_id = %s AND deleted_at IS NULL",
                (document_id,),
            )
            row = cur.fetchone()
            if not row:
                raise HTTPException(404, "Document not found.")
            title = row[1]
            # Soft delete: mark deleted, leave vectors intact (still retrievable via restored copies)
            cur.execute(
                "UPDATE documents SET deleted_at=NOW(), updated_at=NOW(), lifecycle_state='archived' "
                "WHERE document_id=%s",
                (document_id,),
            )
        conn.commit()
    finally:
        conn.close()
    await audit(user.user_id, "DOCUMENT_DELETE", resource_id=document_id,
                details={"title": title, "username": user.username})


@app.post("/api/documents/{document_id}/cancel-processing")
async def cancel_processing(document_id: str, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager or administrator role.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT status FROM documents WHERE document_id=%s AND deleted_at IS NULL",
                (document_id,),
            )
            row = cur.fetchone()
            if not row:
                raise HTTPException(404, "Document not found.")
            if row[0] != "processing":
                raise HTTPException(400, f"Document is not currently processing (status: {row[0]}).")
            cur.execute(
                "UPDATE documents SET status='failed', processing_stage=NULL, updated_at=NOW() "
                "WHERE document_id=%s",
                (document_id,),
            )
        conn.commit()
    finally:
        conn.close()
    await audit(user.user_id, "PROCESSING_CANCELLED", resource_id=document_id,
                details={"username": user.username})
    return {"status": "cancelled", "document_id": document_id}


@app.post("/api/documents/reindex", status_code=202)
async def reindex(body: dict, background_tasks: BackgroundTasks, user: UserContext = Depends(get_user_context)):
    if not user.can_manage_documents():
        raise HTTPException(403, "Requires knowledge-manager, administrator, or system-administrator role.")
    document_id = body.get("document_id")
    if not document_id:
        raise HTTPException(400, "document_id required. Full corpus reindex coming in Phase 6.")
    conn = get_conn()
    try:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT file_path, file_type, title, owner, department, classification "
                "FROM documents WHERE document_id = %s AND status IN ('completed','failed','pending')",
                (document_id,),
            )
            row = cur.fetchone()
    finally:
        conn.close()
    if not row:
        raise HTTPException(404, "Document not found or already processing.")
    file_path, file_type, title, owner, department, classification = row
    metadata = {"title": title, "owner": owner, "department": department, "classification": classification}
    background_tasks.add_task(process_document, document_id, Path(file_path), file_type, metadata)
    return {"status": "queued", "document_id": document_id}
